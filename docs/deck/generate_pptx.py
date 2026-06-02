#!/usr/bin/env python3
"""
Generate agentic-deck.pptx from the HTML deck content.
Run:   python3 docs/deck/generate_pptx.py
Output: docs/deck/agentic-deck.pptx
Import into Google Slides via File → Import slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0F, 0x11, 0x17)
SURFACE = RGBColor(0x16, 0x1B, 0x22)
SURFACE2= RGBColor(0x1C, 0x21, 0x28)
BORDER  = RGBColor(0x30, 0x36, 0x3D)
TEXT    = RGBColor(0xE6, 0xED, 0xF3)
MUTED   = RGBColor(0x8B, 0x94, 0x9E)
ACCENT  = RGBColor(0x00, 0xC4, 0x8C)
ACCENT_BG = RGBColor(0x00, 0x1A, 0x12)
BLUE    = RGBColor(0x79, 0xC0, 0xFF)
BLUE_BG = RGBColor(0x07, 0x14, 0x1F)
BAD     = RGBColor(0xFF, 0x7B, 0x7B)
BAD_BG  = RGBColor(0x1F, 0x07, 0x07)
PURPLE  = RGBColor(0xD2, 0xA8, 0xFF)
ORANGE  = RGBColor(0xFF, 0xA6, 0x57)

# ── Slide dimensions (16:9 widescreen 13.33" × 7.5") ─────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)
ML = Inches(1.0)   # left margin
MT = Inches(0.55)  # top margin
CW = Inches(11.33) # content width

MSO_RECT = 1  # MSO_SHAPE_TYPE.RECTANGLE


# ═══════════════════════════════════════════════════════════════════════════════
# CORE HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def txt(s, text, x, y, w, h, size=13, bold=False, italic=False,
        color=TEXT, align=PP_ALIGN.LEFT, mono=False):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if mono:
        run.font.name = "Courier New"
    return tb


def multi(s, lines, x, y, w, h, size=12, color=TEXT, bold=False,
          align=PP_ALIGN.LEFT, spacing=None):
    """Add textbox with multiple paragraphs. lines = list of str or dict."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.space_before = Pt(spacing)
        if isinstance(line, dict):
            run = p.add_run()
            run.text = line.get("t", "")
            run.font.size = Pt(line.get("sz", size))
            run.font.bold = line.get("b", bold)
            run.font.color.rgb = line.get("c", color)
            run.font.italic = line.get("i", False)
        else:
            run = p.add_run()
            run.text = str(line)
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def lbl(s, text, y=None):
    """Small uppercase accent label."""
    y = y if y is not None else MT
    txt(s, text.upper(), ML, y, CW, Inches(0.28), size=9, bold=True, color=ACCENT)


def hl(s, text, y=None, size=22):
    """Slide headline."""
    y = y if y is not None else MT + Inches(0.35)
    txt(s, text, ML, y, CW, Inches(1.2), size=size, bold=True, color=TEXT)


def divider(s, y):
    sh = s.shapes.add_shape(MSO_RECT, ML, y, Inches(0.45), Pt(3))
    sh.fill.solid()
    sh.fill.fore_color.rgb = ACCENT
    sh.line.fill.background()


def box(s, x, y, w, h, fill=SURFACE, border=BORDER):
    sh = s.shapes.add_shape(MSO_RECT, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def caption(s, text, y):
    box(s, ML, y, CW, Inches(0.55), SURFACE, BORDER)
    txt(s, text, ML + Inches(0.15), y + Inches(0.08),
        CW - Inches(0.3), Inches(0.45), size=10.5, color=MUTED)


def tbl(s, x, y, w, h, rows_data, headers=None,
        header_fill=SURFACE2, row_fill=SURFACE,
        col_widths=None, col0_color=TEXT, col1_color=MUTED,
        header_color=MUTED, size=11):
    """Render a dark-themed table. rows_data = list of tuples."""
    nrows = len(rows_data) + (1 if headers else 0)
    ncols = len(rows_data[0])
    shape = s.shapes.add_table(nrows, ncols, x, y, w, h)
    table = shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = cw

    def style(cell, text, c=TEXT, fill=SURFACE, sz=size, b=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.color.rgb = c

    row_offset = 0
    if headers:
        for i, h_text in enumerate(headers):
            style(table.cell(0, i), h_text, c=header_color,
                  fill=header_fill, sz=9, b=True)
        row_offset = 1

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            c = col0_color if ci == 0 else (ACCENT if ci == 1 and col1_color == ACCENT else col1_color)
            style(table.cell(ri + row_offset, ci), cell_text, c=c, fill=row_fill)

    return shape


def card(s, x, y, w, h, label_text, label_color, title, body,
         fill=SURFACE, border=BORDER):
    """Render a card with label, title, body."""
    box(s, x, y, w, h, fill, border)
    txt(s, label_text.upper(), x + Inches(0.15), y + Inches(0.12),
        w - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=label_color)
    if title:
        txt(s, title, x + Inches(0.15), y + Inches(0.38),
            w - Inches(0.25), Inches(0.45), size=13, bold=True, color=TEXT)
    if body:
        txt(s, body, x + Inches(0.15), y + Inches(0.85),
            w - Inches(0.25), h - Inches(1.0), size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def s01_cover(prs):
    s = blank(prs)
    lbl(s, "software-dev-agentic · April 2026")
    txt(s, "We Gave Our Engineers an", ML, MT + Inches(0.4), CW, Inches(0.9),
        size=40, bold=True, color=TEXT)
    txt(s, "AI Team", ML, MT + Inches(1.3), Inches(5), Inches(0.9),
        size=40, bold=True, color=ACCENT)
    divider(s, MT + Inches(2.35))
    txt(s, "A shared AI engineering toolkit that builds features consistently across every platform — "
        "from a single source of truth, with no per-platform agent changes needed.",
        ML, MT + Inches(2.55), Inches(7.5), Inches(0.85), size=14, color=MUTED)
    tags = ["iOS · Swift/UIKit", "Flutter · Dart/BLoC", "Web · Next.js 15", "Android · Kotlin/MVP"]
    tx = ML
    for tag in tags:
        bw = Inches(2.65)
        box(s, tx, MT + Inches(3.6), bw, Inches(0.42), ACCENT_BG, ACCENT)
        txt(s, tag, tx + Inches(0.14), MT + Inches(3.67), bw - Inches(0.25),
            Inches(0.32), size=11, bold=True, color=ACCENT)
        tx += Inches(2.8)


def s02_background(prs):
    s = blank(prs)
    lbl(s, "Context")
    hl(s, "Software engineering was never about coding. AI just made it obvious.", size=20)
    ty = MT + Inches(1.5)
    tbl(s, ML, ty, CW, Inches(3.6),
        [
            ("Coding was the slowest, largest, and most expensive part of building software.",
             "AI writes code. Coding is no longer the bottleneck."),
            ("Engineers spent most of their time typing — so software engineering looked like typing.",
             "The bottleneck shifted to the idea: what system to build, how it scales, how it holds under load."),
            ('"Software engineering = coding." That assumption made sense. Coding was the constraint.',
             "Software engineering = building stable, scalable systems aligned with product needs. It always was."),
        ],
        headers=["Before AI", "With AI"],
        col0_color=MUTED, col1_color=ACCENT,
        col_widths=[Inches(5.5), Inches(5.83)])
    caption(s, "AI didn't change what software engineering is. It confirmed what it always was. "
            "The engineer's job is now fully about the system — the idea, the architecture, the decisions.",
            MT + Inches(5.3))


def s03_shift(prs):
    s = blank(prs)
    lbl(s, "The Shift")
    hl(s, "Every discipline had its own 'coding'. AI just removed all of them.", size=20)
    ty = MT + Inches(1.5)
    tbl(s, ML, ty, CW, Inches(3.8),
        [
            ("Requirement", "Writing specs, PRDs, user stories",
             "Deciding what problem is worth solving — and what to leave out"),
            ("Design", "Creating mockups, iterating on pixels, building prototypes",
             "How users experience the product. The flow. The information architecture."),
            ("Implementation", "Writing code and boilerplate",
             "Architecture. What to build and how it scales and holds up."),
            ("Testing", "Writing test cases and running manual checks",
             "What 'quality' means. What failure looks like. What's acceptable."),
            ("Delivery", "Configuring pipelines and infrastructure",
             "How reliably and safely the product reaches users."),
        ],
        headers=["Process", "The old bottleneck", "The system — what remains"],
        col0_color=TEXT, col1_color=MUTED, col_widths=[Inches(1.5), Inches(4.5), Inches(5.33)])
    caption(s, "AI executes the bottleneck in every process. What remains is the judgment — "
            "the decisions only your team can make.", MT + Inches(5.5))


def s04_goal(prs):
    s = blank(prs)
    lbl(s, "The Goal")
    hl(s, "Specialized agents for every process. So your team owns the system, not the execution.",
       size=18)
    phases = [
        ("Requirement", "Agent handles", "Drafting specs, PRDs, user stories, acceptance criteria",
         "You focus on", "What problem is worth solving. What to leave out."),
        ("Design", "Agent handles", "Generating mockups, component variants, design tokens",
         "You focus on", "User flow. Information architecture. Experience decisions."),
        ("Implementation", "Agent handles", "Writing code, boilerplate, layer wiring, conventions",
         "You focus on", "Architecture. What to build and how it scales.", True),
        ("Testing", "Agent handles", "Writing test cases, running regression, coverage reports",
         "You focus on", "What quality means. What failure looks like."),
        ("Delivery", "Agent handles", "Pipeline config, infra setup, deployment scripts",
         "You focus on", "Reliability. Rollback strategy. How safely it ships."),
    ]
    card_w = Inches(2.15)
    gap = Inches(0.1)
    ty = MT + Inches(1.45)
    th = Inches(3.8)
    for i, ph in enumerate(phases):
        cx = ML + i * (card_w + gap)
        highlighted = len(ph) > 5 and ph[5]
        fill = ACCENT_BG if highlighted else SURFACE
        border = ACCENT if highlighted else BORDER
        box(s, cx, ty, card_w, th, fill, border)
        txt(s, ph[0].upper(), cx + Inches(0.12), ty + Inches(0.12),
            card_w - Inches(0.2), Inches(0.22), size=8.5, bold=True, color=ACCENT)
        txt(s, ph[1].upper(), cx + Inches(0.12), ty + Inches(0.42),
            card_w - Inches(0.2), Inches(0.18), size=7.5, bold=True, color=MUTED)
        txt(s, ph[2], cx + Inches(0.12), ty + Inches(0.64),
            card_w - Inches(0.2), Inches(1.0), size=10, color=MUTED)
        txt(s, ph[3].upper(), cx + Inches(0.12), ty + Inches(1.75),
            card_w - Inches(0.2), Inches(0.18), size=7.5, bold=True, color=ACCENT)
        txt(s, ph[4], cx + Inches(0.12), ty + Inches(1.97),
            card_w - Inches(0.2), Inches(1.0), size=10, color=TEXT)
    caption(s, "Specialized beats generic. A general-purpose AI assistant doesn't know your architecture, "
            "your conventions, or your layer contracts. These agents do — purpose-built per discipline.",
            MT + Inches(5.45))


def s05_vision(prs):
    s = blank(prs)
    lbl(s, "The Vision")
    hl(s, "The end goal. One prompt. Requirement to shipped.", size=20)
    steps = [
        ("you describe what to build", ACCENT, None),
        ("Requirement", ACCENT, "Spec · priorities · acceptance criteria"),
        ("Design",      ACCENT, "Flows · components · design decisions"),
        ("Implementation", ACCENT, "Code · architecture · all layers, all platforms  [today]"),
        ("Testing",    ACCENT, "Test cases · coverage · regression verified"),
        ("Delivery",   ACCENT, "Pipeline · infra · deployment  [shipped ✓]"),
    ]
    y = MT + Inches(1.5)
    for i, (name, col, detail) in enumerate(steps):
        if i == 0:
            box(s, ML, y, Inches(3.5), Inches(0.38), ACCENT_BG, ACCENT)
            txt(s, name, ML + Inches(0.14), y + Inches(0.07),
                Inches(3.2), Inches(0.28), size=11, bold=True, color=ACCENT)
        else:
            box(s, ML, y, Inches(1.9), Inches(0.38),
                ACCENT_BG if name == "Implementation" else SURFACE,
                ACCENT if name == "Implementation" else BORDER)
            txt(s, name.upper(), ML + Inches(0.12), y + Inches(0.09),
                Inches(1.7), Inches(0.26), size=9, bold=True, color=col)
            if detail:
                txt(s, "→  " + detail, ML + Inches(2.1), y + Inches(0.08),
                    Inches(8.0), Inches(0.28), size=11, color=MUTED)
            if i < len(steps) - 1:
                approve = "you approve" if name != "Delivery" else ""
                if approve:
                    txt(s, approve.upper(), ML + Inches(10.0), y + Inches(0.09),
                        Inches(1.33), Inches(0.26), size=8, bold=True, color=MUTED)
        y += Inches(0.52)
        if i < len(steps) - 1:
            txt(s, "↓", ML + Inches(0.8), y - Inches(0.22),
                Inches(0.5), Inches(0.28), size=13, color=MUTED)
    caption(s, "Not five separate tools. One connected system where each agent's output feeds the next. "
            "The human approves at each gate — the agents handle everything in between.",
            H - Inches(0.85))


def s06_roadmap(prs):
    s = blank(prs)
    lbl(s, "Roadmap")
    hl(s, "A continuous cycle. Core loop + three parallel research tracks.", size=20)

    col_w = Inches(5.4)
    ty = MT + Inches(1.45)

    txt(s, "CORE LOOP — REPEATING CYCLE", ML, ty - Inches(0.28),
        col_w, Inches(0.22), size=8.5, bold=True, color=ACCENT)
    core = [
        ("01", "Foundation",
         "Validate the anatomy. Does it hold across AI platforms and real codebases?",
         "Output: Design Principles — stable agent and skill design, ready to build against."),
        ("02", "Build",
         "Construct agents, personas, and orchestrator skills against the stable design.",
         "Output: Implementation — agents, skills, and personas shipped into lib/."),
        ("03", "Release",
         "Package and distribute a versioned, stable toolkit to all consuming projects.",
         "Output: Distribution — plugin, submodule, or symlink install across every platform."),
        ("04", "Evaluation",
         "Measure effectiveness from real usage. Identify gaps, regressions, or invalidated assumptions.",
         "Output: Report — performance findings against measurable dimensions."),
    ]
    item_h = Inches(1.05)
    for i, (num, title, body, output) in enumerate(core):
        iy = ty + i * item_h
        txt(s, num, ML, iy, Inches(0.4), Inches(0.9), size=12, bold=True, color=ACCENT, mono=True)
        txt(s, title, ML + Inches(0.45), iy, col_w - Inches(0.45), Inches(0.35),
            size=13, bold=True, color=TEXT)
        txt(s, body, ML + Inches(0.45), iy + Inches(0.32), col_w - Inches(0.45),
            Inches(0.38), size=10.5, color=MUTED)
        txt(s, output, ML + Inches(0.45), iy + Inches(0.68), col_w - Inches(0.45),
            Inches(0.33), size=10, color=ACCENT)

    rx = ML + col_w + Inches(0.5)
    txt(s, "RESEARCH TRACKS — PARALLEL", rx, ty - Inches(0.28),
        col_w, Inches(0.22), size=8.5, bold=True, color=MUTED)
    research = [
        ("A", "Collaboration", "research & prototyping",
         "How do multiple teams contribute without breaking each other? "
         "Contribution workflow, ownership model, versioning, and review process."),
        ("B", "Distribution", "research & prototyping",
         "How does the toolkit reach every project reliably? "
         "Submodule + symlinks vs. monorepo vs. plugin — each with different trade-offs."),
        ("C", "Knowledge Architecture", "backlog",
         "How does the knowledge base scale as agents and platforms grow? "
         "Reference doc structure, authoring conventions, versioning, and retrieval."),
        ("D", "Expansion", "backlog",
         "How do we extend specialized agents to every SDLC process? "
         "From Implementation-only today to Requirement, Design, Testing, and Delivery."),
    ]
    for i, (num, title, status, body) in enumerate(research):
        iy = ty + i * item_h
        txt(s, num, rx, iy, Inches(0.4), Inches(0.9), size=12, bold=True, color=MUTED, mono=True)
        txt(s, title + "  [" + status + "]",
            rx + Inches(0.45), iy, col_w - Inches(0.45), Inches(0.35),
            size=13, bold=True, color=TEXT)
        txt(s, body, rx + Inches(0.45), iy + Inches(0.38), col_w - Inches(0.45),
            Inches(0.62), size=10.5, color=MUTED)

    caption(s, "Foundation → Build repeats as a cycle. Research tracks run in parallel and feed findings back continuously.",
            H - Inches(0.85))


def s07_anatomy(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Anatomy")
    hl(s, "Five components. Each with one clear responsibility.", size=20)

    components = [
        ("Reference", TEXT, "The Knowledge",
         "Persistent instruction files Claude loads into its context. CLAUDE.md is always loaded "
         "project-wide. Any other file can be read on demand — bringing targeted knowledge "
         "into the session when needed."),
        ("Skill", ACCENT, "The Hands",
         "Procedural instructions that run in the caller's active session. A named SKILL.md file — "
         "Claude reads it and executes the steps as part of the conversation. Has full access to "
         "Claude tools, including spawning agents and calling other skills."),
        ("Agent", BLUE, "The Brain",
         "A specialized AI assistant spawned in its own isolated context window. Has its own system "
         "prompt, tool access, and permissions — fully independent from the session that spawned it. "
         "Can call skills and spawn further agents. Supports multiple named modes. Returns its result "
         "to the caller when done."),
        ("MCP", PURPLE, "The Reach",
         "Structured tool calls that reach external systems — no file reads needed. Jira tickets, "
         "IDE diagnostics, Figma designs, and build results returned directly to the agent. "
         "Examples: Atlassian (Jira/Confluence/Bitbucket), ide (diagnostics), "
         "Figma, XcodeBuildMCP."),
        ("Hooks", ORANGE, "The Automation",
         "Shell commands that fire on lifecycle events with no model involvement. Logic that must "
         "always run — regardless of what the agent decides to do. Four events: PreToolUse, "
         "PostToolUse, Stop, Notification. For validation and safety guards — not logic that "
         "belongs in a skill."),
    ]
    card_w = Inches(2.15)
    gap    = Inches(0.1)
    ty     = MT + Inches(1.45)
    th     = Inches(4.5)

    for i, (name, col, subtitle, body) in enumerate(components):
        cx = ML + i * (card_w + gap)
        box(s, cx, ty, card_w, th, SURFACE, BORDER)
        txt(s, name.upper(), cx + Inches(0.12), ty + Inches(0.12),
            card_w - Inches(0.2), Inches(0.22), size=8.5, bold=True, color=col)
        txt(s, subtitle, cx + Inches(0.12), ty + Inches(0.38),
            card_w - Inches(0.2), Inches(0.35), size=11, bold=True, color=TEXT)
        txt(s, body, cx + Inches(0.12), ty + Inches(0.78),
            card_w - Inches(0.2), th - Inches(0.9), size=10, color=MUTED)


def s07b_capability_matrix(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Capabilities")
    hl(s, "What each component is, by nature.", size=22)
    ty = MT + Inches(1.4)

    capabilities = [
        "LLM reasoning",
        "Isolated context window",
        "Runs in caller's context",
        "Uses Claude tools",
        "Writes source files",
        "Multiple invocation modes",
        "Spawns agents",
        "Calls skills",
        "Bridges external systems",
        "Grep-addressable knowledge",
        "Shell execution (no model)",
    ]
    checkmarks = [
        # Reference, Skill, Agent, MCP, Hooks
        ("—", "—", "✓", "—", "—"),
        ("—", "—", "✓", "—", "—"),
        ("—", "✓", "—", "—", "—"),
        ("—", "✓", "✓", "—", "—"),
        ("—", "✓", "✓", "—", "—"),
        ("—", "—", "✓", "—", "—"),
        ("—", "✓", "✓", "—", "—"),
        ("—", "✓", "✓", "—", "—"),
        ("—", "—", "—", "✓", "—"),
        ("✓", "—", "—", "—", "—"),
        ("—", "—", "—", "—", "✓"),
    ]
    col_colors = [TEXT, MUTED, ACCENT, BLUE, PURPLE, ORANGE]
    headers = ["Capability", "Reference", "Skill", "Agent", "MCP", "Hooks"]
    col_widths = [Inches(3.5), Inches(1.55), Inches(1.55), Inches(1.55), Inches(1.55), Inches(1.58)]

    nrows = len(capabilities) + 1
    ncols = 6
    shape = s.shapes.add_table(nrows, ncols, ML, ty, CW, Inches(5.3))
    table = shape.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = cw

    def st(cell, t, c=TEXT, fill=SURFACE, sz=11, b=False, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c

    for ci, h in enumerate(headers):
        st(table.cell(0, ci), h, c=col_colors[ci] if ci > 0 else MUTED,
           fill=SURFACE2, sz=9, b=True,
           align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    for ri, (cap, checks) in enumerate(zip(capabilities, checkmarks)):
        st(table.cell(ri+1, 0), cap, c=TEXT, fill=SURFACE)
        for ci, mark in enumerate(checks):
            fill = SURFACE
            c = ACCENT if mark == "✓" else MUTED
            st(table.cell(ri+1, ci+1), mark, c=c, fill=fill,
               align=PP_ALIGN.CENTER, b=(mark == "✓"))


def s_persona(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Persona")
    hl(s, "A Persona mirrors a role. An Orchestrator skill mirrors what that role does.", size=18)

    ty = MT + Inches(1.45)
    lw = Inches(4.8)
    rw = Inches(6.0)
    rx = ML + lw + Inches(0.53)

    txt(s, "PERSONA = SDLC PHASE + THE ROLE THAT OWNS IT",
        ML, ty - Inches(0.25), lw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    tbl(s, ML, ty, lw, Inches(3.2),
        [
            ("Implementation", "Software Engineer", "developer"),
            ("Testing",        "QA Engineer",       "qa"),
            ("Requirement",    "Product Manager",   "— research"),
            ("Design",         "Designer",          "— research"),
            ("Delivery",       "DevOps / EM",       "— research"),
        ],
        headers=["SDLC Phase", "Role", "Persona"],
        col0_color=TEXT, col1_color=MUTED, col_widths=[Inches(1.6), Inches(1.7), Inches(1.5)])

    txt(s, "ORCHESTRATOR SKILL = THE ROLE'S ACTUAL WORKFLOW",
        rx, ty - Inches(0.25), rw, Inches(0.22), size=8.5, bold=True, color=MUTED)

    personas = [
        ("developer persona", ACCENT,
         [("/developer-rfc",
           "Break down a Jira Epic + PRD into RFC + ticket breakdown"),
          ("/developer-plan-feature",
           "Plan and build a feature from a ticket or RFC — convergence loop → implementation")]),
        ("qa persona", BLUE,
         [("/qa-generate-testcase",
           "Break down PRD, RFC, or Figma into test cases — outputs CSV + Jira comment")]),
    ]
    py = ty
    for (pname, pcol, skills_list) in personas:
        ph = Inches(1.35) if len(skills_list) > 1 else Inches(0.9)
        box(s, rx, py, rw, ph, SURFACE2, BORDER)
        txt(s, pname.upper(), rx + Inches(0.15), py + Inches(0.1),
            rw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=pcol)
        for si, (cmd, desc) in enumerate(skills_list):
            sy = py + Inches(0.38) + si * Inches(0.42)
            txt(s, cmd, rx + Inches(0.15), sy, Inches(2.4), Inches(0.35),
                size=10, bold=True, color=pcol, mono=True)
            txt(s, desc, rx + Inches(2.65), sy, rw - Inches(2.8), Inches(0.35),
                size=10, color=MUTED)
        py += ph + Inches(0.18)


def s_example(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · In Practice")
    hl(s, "One request. The whole system in motion.", size=22)

    flow = [
        ("engineer", "/developer-plan-feature \"OTP Login\"", ACCENT, None),
        ("Skill", "developer-plan-feature",
         ACCENT, "Owns the session — routes the request, manages approval gates, loops until converged"),
        ("Agents", "strategist · 3× planners",
         BLUE, "Each in its own isolated context — strategist decides scope, "
               "planners reason per layer (domain · data · presentation) in parallel"),
        ("Knowledge", "reference/",
         MUTED, "Agents load only the pattern sections their layer needs — "
                "domain contracts, data models, UI conventions. Nothing more."),
        ("Output", "plan → code",
         ACCENT, "Plan reviewed and approved by engineer → feature worker writes all layers "
                 "across every platform  [shipped ✓]"),
    ]
    y = MT + Inches(1.45)
    for i, row in enumerate(flow):
        kind, name, col, detail = row
        bw = Inches(2.3)
        fill = ACCENT_BG if col == ACCENT else (BLUE_BG if col == BLUE else SURFACE)
        border = ACCENT if col == ACCENT else (BLUE if col == BLUE else BORDER)
        box(s, ML, y, bw, Inches(0.48), fill, border)
        txt(s, kind.upper(), ML + Inches(0.12), y + Inches(0.06),
            bw - Inches(0.2), Inches(0.2), size=8, bold=True, color=col)
        txt(s, name, ML + Inches(0.12), y + Inches(0.26),
            bw - Inches(0.2), Inches(0.22), size=10, color=TEXT, mono=True)
        if detail:
            txt(s, detail, ML + Inches(2.5), y + Inches(0.1),
                Inches(8.83), Inches(0.35), size=11, color=MUTED)
        if i < len(flow) - 1:
            arrow = "↓ spawns" if i == 0 else ("↓ grep" if i == 2 else "↓")
            txt(s, arrow, ML + Inches(0.9), y + Inches(0.55),
                Inches(1.5), Inches(0.28), size=11, color=MUTED)
        y += Inches(0.95)
    caption(s, "The engineer describes what to build. Skill, agents, and knowledge each do exactly one job — nothing more.",
            H - Inches(0.85))


def s_why_not_e2e(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Limitations")
    hl(s, "Why not end-to-end? Four reasons human review still matters.", size=20)

    reasons = [
        ("01", "Context window",
         "A full end-to-end workflow — across multiple roles and hand-offs — fills the main context. "
         "Once compacted, the session reasons on a lossy summary. Quality at the final step is lower "
         "than at the first."),
        ("02", "Token billing",
         "A billing cap or timeout mid-workflow orphans the Orchestrator skill. There is no clean "
         "auto-recovery — a human has to re-enter, re-read state, and decide how to continue. "
         "The longer the chain, the higher the risk."),
        ("03", "Supervision",
         "Agents aren't mature enough to run unsupervised across team boundaries. Each hand-off is a "
         "decision point — does the RFC accurately reflect the Epic? Does the plan match the RFC? "
         "A human eye at each gate catches drift early."),
        ("04", "Gaps",
         "There is no agreed input/output contract between Product → Design → Developer → QA yet. "
         "We're still mapping what each team produces and what the next team needs. "
         "Agents can't bridge a gap that hasn't been defined."),
    ]
    cw = Inches(2.75)
    gap = Inches(0.1)
    ty = MT + Inches(1.5)
    th = Inches(3.8)
    for i, (num, title, body) in enumerate(reasons):
        cx = ML + i * (cw + gap)
        box(s, cx, ty, cw, th, BAD_BG, BAD)
        txt(s, num, cx + Inches(0.15), ty + Inches(0.12),
            cw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=BAD)
        txt(s, title, cx + Inches(0.15), ty + Inches(0.38),
            cw - Inches(0.25), Inches(0.42), size=13, bold=True, color=TEXT)
        txt(s, body, cx + Inches(0.15), ty + Inches(0.88),
            cw - Inches(0.25), th - Inches(1.0), size=11, color=MUTED)
    caption(s, "Each limitation is a research problem, not a dead end. Context efficiency, resume protocols, "
            "maturity metrics, and cross-team I/O contracts are all active work.", H - Inches(0.85))


def s_evolution(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Evolution")
    hl(s, "We didn't design it all at once. Each problem taught us the next piece.", size=18)

    evol = [
        ("Origin", ACCENT, "Raw Skill",
         "SKILL.md in the main context. Engineer types, skill runs. No agents, no isolation. "
         "Worked fine for simple, bounded tasks."),
        ("Problem 1", BAD, "Context fills",
         "As the skill grows, the main context fills. Quality degrades. "
         "Solution: Split to Agents — isolated context windows. Skill orchestrates; agent reasons. "
         "Only the result returns to main context."),
        ("Problem 2", BAD, "One pass isn't enough",
         "Domain, data, and presentation layers depend on each other. A single agent pass misses "
         "cross-layer impact. Solution: Convergence loop. Planners report cross-layer impact. "
         "Strategist decides: more rounds or converged?"),
        ("Problem 3", BAD, "Knowledge Coupling",
         "Embedding knowledge inside agents causes bloat, duplication, and brittle updates. "
         "Solution: Knowledge directory + Grep. Lives in one place, versioned independently. "
         "Agents grep only the section their task needs."),
        ("Problem 4", BAD, "Repeating ourselves",
         "Agents writing the same boilerplate inline, per platform. No reuse. "
         "Solution: Type P skills called by agents — not the Orchestrator. "
         "Platform-specific. One definition, updated once."),
    ]
    cw = Inches(2.15)
    gap = Inches(0.1)
    ty = MT + Inches(1.4)
    th = Inches(2.8)
    for i, (lbl_t, col, title, body) in enumerate(evol):
        cx = ML + i * (cw + gap)
        fill = ACCENT_BG if col == ACCENT else SURFACE
        border = ACCENT if col == ACCENT else BORDER
        box(s, cx, ty, cw, th, fill, border)
        txt(s, lbl_t.upper(), cx + Inches(0.12), ty + Inches(0.1),
            cw - Inches(0.2), Inches(0.2), size=8, bold=True, color=col)
        txt(s, title, cx + Inches(0.12), ty + Inches(0.34),
            cw - Inches(0.2), Inches(0.38), size=12, bold=True, color=TEXT)
        txt(s, body, cx + Inches(0.12), ty + Inches(0.78),
            cw - Inches(0.2), th - Inches(0.9), size=10, color=MUTED)

    txt(s, "SO WHEN BUILDING A NEW ORCHESTRATOR — ANSWER THESE FOUR QUESTIONS FIRST",
        ML, ty + th + Inches(0.18), CW, Inches(0.22),
        size=8.5, bold=True, color=ACCENT)

    questions = [
        ("01", "What's the goal?", "Output",
         "Define the structured result before writing any instructions. Declare it first."),
        ("02", "What does it need?", "Input",
         "Declare every required parameter explicitly. Missing input = MISSING INPUT immediately."),
        ("03", "How does it run?", "Agents + Loop",
         "Convergence loop or not? Decision blocks in, findings files on disk out."),
        ("04", "Will 200K hold?", "Context budget",
         "Estimate rounds × output size. If too large, split into a second Orchestrator."),
    ]
    qcw = Inches(2.75)
    qgap = Inches(0.1)
    qy = ty + th + Inches(0.46)
    qh = Inches(1.6)
    for i, (num, title, sub, body) in enumerate(questions):
        cx = ML + i * (qcw + qgap)
        fill = BAD_BG if i == 3 else SURFACE
        border = BAD if i == 3 else BORDER
        box(s, cx, qy, qcw, qh, fill, border)
        txt(s, num, cx + Inches(0.12), qy + Inches(0.1),
            qcw - Inches(0.2), Inches(0.2), size=8.5, bold=True,
            color=BAD if i == 3 else ACCENT)
        txt(s, title, cx + Inches(0.12), qy + Inches(0.34),
            qcw - Inches(0.2), Inches(0.32), size=12, bold=True, color=TEXT)
        txt(s, sub.upper(), cx + Inches(0.12), qy + Inches(0.7),
            qcw - Inches(0.2), Inches(0.18), size=8, bold=True,
            color=BAD if i == 3 else ACCENT)
        txt(s, body, cx + Inches(0.12), qy + Inches(0.92),
            qcw - Inches(0.2), Inches(0.6), size=10, color=MUTED)


def s_reference_hl(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Reference")
    hl(s, "Knowledge lives in one place. Agents stay focused on their job.", size=20)

    cards_data = [
        ("Knowledge Directory", ACCENT, "What the system knows",
         "Architecture patterns, layer contracts, platform conventions. "
         "Centralized in reference/. Versioned and owned independently of any agent or skill.",
         ACCENT_BG, ACCENT),
        ("Agents & Skills", MUTED, "What the system does",
         "Scoped to one concern: thinking, routing, reasoning, orchestrating. "
         "Load only what their task needs — grep a section, not the full file. "
         "Contain no domain knowledge of their own.",
         SURFACE, BORDER),
        ("The Result", MUTED, "Concerns separated",
         "Update a pattern → change one knowledge file, zero agents. "
         "Add a platform → add a new impl file, no agent rewrite. "
         "Debug a wrong output → check knowledge or agent independently.",
         SURFACE, BORDER),
    ]
    cw = Inches(3.6)
    gap = Inches(0.17)
    ty = MT + Inches(1.5)
    th = Inches(3.5)
    for i, (title, col, subtitle, body, fill, border) in enumerate(cards_data):
        cx = ML + i * (cw + gap)
        box(s, cx, ty, cw, th, fill, border)
        txt(s, title.upper(), cx + Inches(0.15), ty + Inches(0.12),
            cw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=col)
        txt(s, subtitle, cx + Inches(0.15), ty + Inches(0.38),
            cw - Inches(0.25), Inches(0.38), size=13, bold=True, color=TEXT)
        multi(s, ["→  " + b for b in body.split(". ") if b],
              cx + Inches(0.15), ty + Inches(0.85),
              cw - Inches(0.25), th - Inches(1.0), size=11, color=MUTED)
    caption(s, "Agents and skills are logic. The knowledge directory is data. "
            "They change for different reasons, at different times, by different people — so they live apart.",
            H - Inches(0.85))


def s_reference_ll(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Reference  [Low-Level]")
    hl(s, "Topic. Term. Three reference types. One grep key.", size=22)

    ty = MT + Inches(1.45)
    lw = Inches(5.3)
    rw = Inches(5.7)
    rx = ML + lw + Inches(0.33)

    tree_lines = [
        "Process   Implementation · Design · Requirement · …",
        "  │",
        "  Topic   domain · data · presentation · design-system · …",
        "  │            one reference file per topic per platform",
        "  │",
        "  Term    UseCase · Entity · Repository · MpButton · …",
        "  │            one ## heading per term, same name everywhere",
        "  │",
        "  ├── -theory.md   what it IS · platform-agnostic",
        "  ├── -impl.md     how to build it in this language",
        "  └── -catalog.md  queryable symbol inventory (### per entry)",
    ]
    box(s, ML, ty, lw, Inches(2.4), SURFACE, BORDER)
    txt(s, "\n".join(tree_lines), ML + Inches(0.15), ty + Inches(0.12),
        lw - Inches(0.25), Inches(2.2), size=10.5, color=TEXT, mono=True)

    box(s, ML, ty + Inches(2.55), lw, Inches(1.3), ACCENT_BG, ACCENT)
    txt(s, "THE TERM IS THE GREP KEY",
        ML + Inches(0.15), ty + Inches(2.68),
        lw - Inches(0.25), Inches(0.22), size=8, bold=True, color=ACCENT)
    txt(s, "One Term = one canonical name = one ## heading, everywhere.\n"
        "Catalog entries use ### per symbol — agents grep the symbol name directly, "
        "read 8 lines. No full-file load, ever.",
        ML + Inches(0.15), ty + Inches(2.94),
        lw - Inches(0.25), Inches(0.85), size=11, color=TEXT)

    code_lines = [
        "domain-theory.md  — WHAT and WHY",
        "## UseCase",
        "A use case encapsulates one business operation.",
        "Invariants: no framework imports, single execute()...",
        "",
        "domain-impl.md  — HOW in Swift",
        "## UseCase   ← same Term, Swift body",
        "final class GetUserUseCase: UseCase {",
        "  func execute(_ params: Params) async throws",
        "",
        "mekari-pixel-flutter-catalog.md  — WHAT EXISTS",
        "## Atoms",
        "",
        "### MpButton",
        "- Category: `atoms/button`",
        "- Description: Mekari Mobile Kit Button.",
        "- Key params: `label`, `icon`, `onPressed`, `size`",
        "- Variants: `MpButton.primary()`, `.secondary()`",
    ]
    box(s, rx, ty, rw, Inches(3.9), SURFACE, BORDER)
    txt(s, "THREE TYPES — THEORY · IMPL · CATALOG",
        rx + Inches(0.15), ty + Inches(0.1),
        rw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=MUTED)
    txt(s, "\n".join(code_lines),
        rx + Inches(0.15), ty + Inches(0.38),
        rw - Inches(0.25), Inches(3.4), size=10, color=TEXT, mono=True)

    caption(s, "-theory answers what and why.  -impl answers how in that language.  "
            "-catalog answers what exists — a queryable inventory of available symbols.",
            H - Inches(0.85))


def s_skills_hl(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Skills")
    hl(s, "Two skill types. Type O is the workflow. Type P is a step inside it.", size=20)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.5)
    th = Inches(3.9)

    box(s, ML, ty, lw, th, SURFACE, BORDER)
    txt(s, "TYPE O — ORCHESTRATOR", ML + Inches(0.15), ty + Inches(0.12),
        lw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=ACCENT)
    txt(s, "The entry point", ML + Inches(0.15), ty + Inches(0.38),
        lw - Inches(0.25), Inches(0.35), size=14, bold=True, color=TEXT)
    multi(s, [
        "Triggered by the engineer. Owns the full session — routing, approval gates, and "
        "coordinating agents. Simple workflows do their own work. Complex ones delegate.",
        "",
        "Example: /developer-plan-feature",
        "",
        "Owns: Routing · approval gates · the loop. The engineer talks to the skill, not the agents.",
        "Delegates: For decisions or code generation, delegate to dedicated agents — each isolated, "
        "each specialized for one job.",
        "Parallel spawn: One skill step can spawn N agents simultaneously — same wall-clock time.",
        "Convergence: The skill loops — spawns agents, collects decisions, re-spawns — until converged.",
        "Constraint: Runs in your session — every agent result lands here. Context fills. "
        "Design to keep rounds low.",
    ], ML + Inches(0.15), ty + Inches(0.82), lw - Inches(0.25), th - Inches(0.95),
       size=11, color=MUTED)

    box(s, rx, ty, rw, th, SURFACE, BORDER)
    txt(s, "TYPE P — PROCEDURE", rx + Inches(0.15), ty + Inches(0.12),
        rw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=MUTED)
    txt(s, "A thin step", rx + Inches(0.15), ty + Inches(0.38),
        rw - Inches(0.25), Inches(0.35), size=14, bold=True, color=TEXT)
    multi(s, [
        "Called by agents, never by users. Creates one artifact using the right platform conventions. "
        "No decisions. No branching. Just execution.",
        "",
        "Example: developer-domain-create-entity",
        "",
        "→ user-invocable: false",
        "→ No AskUserQuestion — no branching logic",
        "→ Grep-first reference loading (never full file)",
        "→ Check preconditions, STOP if file already exists",
        "→ Logic belongs in the worker — not here",
    ], rx + Inches(0.15), ty + Inches(0.82), rw - Inches(0.25), th - Inches(0.95),
       size=11, color=MUTED)

    caption(s, "A skill makes Claude's behavior consistent and repeatable — same workflow, every engineer, every platform. "
            "Without a Type O, there is no guardrailed entry, no approval gate. Just a raw AI session.",
            H - Inches(0.85))


def s_skill_flow(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Orchestrator Skill  [Low-Level]")
    hl(s, "The Orchestrator skill runs in your context window. Three things it can do — and one constraint.",
       size=18)

    panels = [
        ("Capability 1 — Parallel agents", ACCENT, ACCENT_BG, ACCENT,
         "One skill step can spawn N agents simultaneously. Each runs in its own isolated context "
         "and returns findings. Same wall-clock time as spawning one.",
         "# spawn domain + data + pres planners\n"
         "# in the same skill step\n"
         "Agent(domain-planner, background: true)\n"
         "Agent(data-planner,   background: true)\n"
         "Agent(pres-planner,   background: true)\n"
         "← all three run in parallel"),
        ("Capability 2 — Convergence loop", ACCENT, ACCENT_BG, ACCENT,
         "The skill owns the loop. It spawns agents, collects decisions, and loops back until the "
         "strategist signals convergence. Each agent call is self-contained.",
         "while not converged:\n"
         "  result = Agent(strategist, findings)\n"
         "  # Decision: spawn-planners\n"
         "  findings = spawn_planners_parallel()\n"
         "  # Decision: converged\n"
         "  → break"),
        ("Constraint — Main context window", BAD, BAD_BG, BAD,
         "Every agent result returns to the skill — which lives in your main session. Each round adds "
         "more history. When context fills, Claude compacts it into a summary. "
         "Compaction = reduced quality. Design for minimal rounds. Write findings to disk — "
         "pass paths, not content.",
         ""),
    ]
    cw = Inches(3.6)
    gap = Inches(0.17)
    ty = MT + Inches(1.5)
    th = Inches(4.5)
    for i, (title, col, fill, border, desc, code) in enumerate(panels):
        cx = ML + i * (cw + gap)
        box(s, cx, ty, cw, th, fill, border)
        txt(s, title.upper(), cx + Inches(0.15), ty + Inches(0.12),
            cw - Inches(0.25), Inches(0.22), size=8, bold=True, color=col)
        txt(s, desc, cx + Inches(0.15), ty + Inches(0.4),
            cw - Inches(0.25), Inches(1.8), size=11, color=MUTED)
        if code:
            box(s, cx + Inches(0.1), ty + Inches(2.35), cw - Inches(0.2),
                Inches(1.65), SURFACE, BORDER)
            txt(s, code, cx + Inches(0.2), ty + Inches(2.45),
                cw - Inches(0.35), Inches(1.5), size=10, color=TEXT, mono=True)

    caption(s, "Agents run in isolated contexts — they never see the main session. "
            "But the skill does, and every call back to it has a cost. "
            "Parallel spawning and disk-based findings are how we keep the round count low.",
            H - Inches(0.85))


def s_skills_ll(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Skills  [Low-Level]")
    hl(s, "Two types. Each with a different job, a different caller.", size=22)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.45)
    th = Inches(4.7)

    box(s, ML, ty, lw, th, SURFACE, BORDER)
    txt(s, "O · ORCHESTRATOR — USER ENTRY, OWNS THE SESSION",
        ML + Inches(0.15), ty + Inches(0.1),
        lw - Inches(0.25), Inches(0.22), size=8, bold=True, color=ACCENT)
    code_o = ("--- developer-plan-feature ---\n"
              "user-invocable: true\n"
              "allowed-tools: Agent, AskUserQuestion, Read\n\n"
              "# Step 1 — Gather Intent\n"
              "Spawn strategist (gather-intent)\n"
              "→ returns Decision: spawn-planners\n\n"
              "# Step 2 — Convergence Loop\n"
              "Spawn planners in parallel per round\n"
              "Send findings → strategist\n"
              "Loop until converged or blocked\n\n"
              "# Step 3 — Synthesize\n"
              "Spawn strategist (synthesize)\n"
              "→ writes plan.md + context.md\n\n"
              "# Step 4 — Approve\n"
              "AskUserQuestion\n"
              "  Approve  → execute\n"
              "  Discuss  → revise → re-approve\n"
              "  Discard  → delete run, stop\n\n"
              "# Step 5 — Execute\n"
              "Spawn feature-worker\n"
              "  plan.md + context.md inlined")
    txt(s, code_o, ML + Inches(0.15), ty + Inches(0.38),
        lw - Inches(0.25), th - Inches(0.5), size=10.5, color=TEXT, mono=True)

    box(s, rx, ty, rw, th, SURFACE, BORDER)
    txt(s, "P · PROCEDURE — AGENT CALLS, THIN AND NON-INTERACTIVE",
        rx + Inches(0.15), ty + Inches(0.1),
        rw - Inches(0.25), Inches(0.22), size=8, bold=True, color=MUTED)
    code_p = ("--- developer-domain-create-entity ---\n"
              "user-invocable: false\n"
              "allowed-tools: Read, Grep, Bash\n"
              "(no AskUserQuestion — no branching logic)\n\n"
              "# Input\n"
              "entity_name, target_path, module\n\n"
              "# Step 1 — Load reference (Grep-first)\n"
              "Grep \"^## Entity\" in domain.md\n"
              "Read offset + limit  ← section only\n"
              "Never the full file\n\n"
              "# Step 2 — Check preconditions\n"
              "Glob target path\n"
              "STOP if file already exists\n\n"
              "# Step 3 — Write\n"
              "Create entity from reference pattern\n"
              "Logic belongs in the worker — not here\n\n"
              "# Output\n"
              "Glob + Grep verified path only")
    txt(s, code_p, rx + Inches(0.15), ty + Inches(0.38),
        rw - Inches(0.25), th - Inches(0.5), size=10.5, color=TEXT, mono=True)


def s_agents_hl(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Agents")
    hl(s, "Five parts. Every agent is built from the same anatomy.", size=22)

    parts = [
        ("01", "Input",
         "What the agent receives to start — mode, feature name, platform, file paths. "
         "Declared explicitly. Missing input = immediate stop.", False),
        ("02", "Knowledge",
         "Reference docs and patterns loaded on demand. Each agent greps only what its task needs — "
         "nothing more. Specialization is a loading decision.", False),
        ("03", "Reasoning",
         "The LLM applies thinking, deciding, and branching to everything it has loaded. This is the "
         "part no fixed script can replace — it handles ambiguity and edge cases.", True),
        ("04", "Output",
         "Declared and structured — a Decision: block, a ## Findings section, verified file paths. "
         "The caller routes on it without ambiguity.", False),
        ("05", "Modes",
         "One agent, many invocation contexts. Each mode loads only its own instruction block. "
         "The agent never reads instructions for modes it isn't in.", False),
    ]
    cw = Inches(2.15)
    gap = Inches(0.1)
    ty = MT + Inches(1.5)
    th = Inches(4.5)
    for i, (num, title, body, highlight) in enumerate(parts):
        cx = ML + i * (cw + gap)
        fill = ACCENT_BG if highlight else SURFACE
        border = ACCENT if highlight else BORDER
        box(s, cx, ty, cw, th, fill, border)
        txt(s, num, cx + Inches(0.15), ty + Inches(0.12),
            cw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=ACCENT)
        txt(s, title, cx + Inches(0.15), ty + Inches(0.38),
            cw - Inches(0.25), Inches(0.38), size=14, bold=True, color=TEXT)
        txt(s, body, cx + Inches(0.15), ty + Inches(0.85),
            cw - Inches(0.25), th - Inches(1.0), size=11, color=MUTED)
    caption(s, "Isolated context window keeps all five parts — and everything they touch — out of your main session. "
            "The main session only ever sees the structured output.", H - Inches(0.85))


def s_agents_why(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Agents")
    hl(s, "Why this anatomy. Agents aren't black boxes — every behavior traces back to one part.",
       size=18)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.5)

    txt(s, "IF THE AGENT BEHAVES WRONG — CHECK WHICH PART",
        ML, ty - Inches(0.28), lw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    tbl(s, ML, ty, lw, Inches(4.5),
        [
            ("Agent starts wrong, misses context, or acts on wrong scope", "Input"),
            ("Agent doesn't know the right patterns or uses wrong conventions", "Knowledge"),
            ("Agent makes wrong decisions or draws wrong conclusions", "Reasoning"),
            ("Agent returns incomplete, unparseable, or unexpected result", "Output"),
            ("Agent behaves inconsistently across different invocations", "Mode"),
        ],
        headers=["Symptom", "Where to look"],
        col0_color=MUTED, col1_color=ACCENT,
        col_widths=[Inches(3.9), Inches(1.6)])

    txt(s, "EACH PART IS INDEPENDENTLY REPLACEABLE",
        rx, ty - Inches(0.28), rw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    replacements = [
        ("Input",
         "Add a new source — Jira, Figma, local file — without changing reasoning or output."),
        ("Knowledge",
         "Swap Grep-loaded docs for vector search, a database, or a catalog — same slot, richer retrieval."),
        ("Reasoning",
         "Today: LLM. Tomorrow: a deterministic rule engine for the same decision slot, "
         "if the logic is well-understood enough."),
        ("Output",
         "Change the schema the caller expects without rewriting the agent's reasoning."),
        ("Mode",
         "Add a new invocation context — a new mode — without affecting existing ones."),
    ]
    ry = ty
    rh = Inches(0.78)
    gap = Inches(0.1)
    for part, desc in replacements:
        box(s, rx, ry, rw, rh, SURFACE, BORDER)
        txt(s, part, rx + Inches(0.15), ry + Inches(0.1),
            Inches(1.0), Inches(0.25), size=11, bold=True, color=ACCENT)
        txt(s, desc, rx + Inches(1.25), ry + Inches(0.1),
            rw - Inches(1.4), Inches(0.6), size=11, color=MUTED)
        ry += rh + gap


def s_agent_anatomy_ll(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · Agents  [Low-Level]")
    hl(s, "One file. Five slots. Each with one clear job.", size=22)

    lw = Inches(5.8)
    rw = Inches(5.2)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.45)
    th = Inches(5.1)

    box(s, ML, ty, lw, th, SURFACE, BORDER)
    txt(s, "DEVELOPER-FEATURE-STRATEGIST (abbreviated)",
        ML + Inches(0.15), ty + Inches(0.1),
        lw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=BLUE)
    code = ("--- frontmatter ---\n"
            "model: claude-sonnet-4-6\n"
            "tools: Read, Glob, Grep\n\n"
            "## Input                        ← (1) Input\n"
            "mode: gather-intent | process-findings | synthesize | resume\n"
            "feature_name: <string>\n"
            "Missing input → MISSING INPUT: <param>  (stop immediately)\n\n"
            "## Knowledge                    ← (2) Knowledge\n"
            "Grep reference/domain-theory.md ^## UseCase\n"
            "One targeted section loaded — never the full file\n\n"
            "## Reasoning                    ← (3) Reasoning\n"
            "Assess existing runs. Decide which layers need planners.\n"
            "Handle ambiguity: overlapping modules, uncertain scope.\n\n"
            "## Output                       ← (4) Output\n"
            "Decision: spawn-planners [domain, data, pres]\n"
            "Structured block — caller routes without parsing\n\n"
            "## Modes                        ← (5) Modes\n"
            "[gather-intent]     load only when mode=gather-intent\n"
            "[process-findings]  load only when mode=process-findings\n"
            "One body — each call loads only its mode's instruction block")
    txt(s, code, ML + Inches(0.15), ty + Inches(0.38),
        lw - Inches(0.25), th - Inches(0.5), size=10, color=TEXT, mono=True)

    slots = [
        ("01 Input", "Declared explicitly. Missing input stops the agent immediately with a clear message — "
                     "no guessing, no partial execution."),
        ("02 Knowledge", "Grep-targeted. One section per task — not the whole file. "
                         "Change what it loads, change what it knows."),
        ("03 Reasoning", "The only part that can't be scripted. Handles ambiguity — "
                         "and the only slot that could one day be a deterministic engine."),
        ("04 Output", "Structured and declared. The calling skill routes on a Decision: block — "
                      "no parsing, no inference."),
        ("05 Modes", "One agent file, many invocation contexts. Each mode is a named section — "
                     "only that block is read per call."),
    ]
    ry = ty
    slot_h = Inches(0.88)
    for i, (slot_name, slot_desc) in enumerate(slots):
        fill = ACCENT_BG if i == 2 else SURFACE
        border = ACCENT if i == 2 else BORDER
        box(s, rx, ry, rw, slot_h, fill, border)
        txt(s, slot_name, rx + Inches(0.15), ry + Inches(0.1),
            Inches(1.2), Inches(0.25), size=11, bold=True, color=ACCENT)
        txt(s, slot_desc, rx + Inches(1.45), ry + Inches(0.1),
            rw - Inches(1.6), Inches(0.7), size=11, color=MUTED)
        ry += slot_h + Inches(0.1)


def s_mcp_hl(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · MCP")
    hl(s, "Agents reach your tools. You stop pasting.", size=22)
    divider(s, MT + Inches(1.45))
    tbl(s, ML, MT + Inches(1.65), CW, Inches(3.5),
        [
            ('"Paste the Jira ticket so I can understand the requirements."',
             "Agent fetches the ticket directly — description, AC, linked issues, comments."),
            ('"Paste the Figma design or take a screenshot of each screen."',
             "Agent pulls design data, component specs, and layout files from Figma directly."),
            ('"Run the build and paste the error so I can see what\'s failing."',
             "Agent triggers the build, reads the result, and acts on it — no relay needed."),
        ],
        headers=["Before MCP", "With MCP"],
        col0_color=MUTED, col1_color=ACCENT,
        col_widths=[Inches(5.5), Inches(5.83)])
    caption(s, "MCP turns external tools into first-class inputs. The agent's context is always fresh and complete — "
            "no copy-paste, no stale screenshots, no information lost in translation.", H - Inches(0.85))


def s_mcp_ll(prs):
    s = blank(prs)
    lbl(s, "Phase 01 — Foundation · Design Principle · MCP  [Low-Level]")
    hl(s, "Skills reach the codebase. MCP reaches everything else.", size=22)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.45)

    txt(s, "TWO DIRECTIONS OF REACH",
        ML, ty - Inches(0.28), lw, Inches(0.22), size=8.5, bold=True, color=MUTED)

    box(s, ML, ty, lw, Inches(1.5), SURFACE, BORDER)
    txt(s, "VIA SKILL — INTO THE CODEBASE",
        ML + Inches(0.15), ty + Inches(0.1),
        lw - Inches(0.25), Inches(0.22), size=8, bold=True, color=ACCENT)
    txt(s, "# Worker calls a platform-contract skill\n"
        "call developer-domain-create-entity\n"
        "  feature: GetUser\n"
        "  platform: ios\n\n"
        "→ creates GetUserUseCase.swift",
        ML + Inches(0.15), ty + Inches(0.36),
        lw - Inches(0.25), Inches(1.0), size=10.5, color=TEXT, mono=True)

    box(s, ML, ty + Inches(1.65), lw, Inches(1.5), SURFACE, BORDER)
    txt(s, "VIA MCP — INTO EXTERNAL SYSTEMS",
        ML + Inches(0.15), ty + Inches(1.75),
        lw - Inches(0.25), Inches(0.22), size=8, bold=True, color=PURPLE)
    txt(s, "# Agent calls an MCP tool directly\n"
        "mcp__atlassian__get_jira_issue(\n"
        '  issue_key: "PROJ-42"\n'
        ")\n\n"
        "→ returns ticket, context, acceptance criteria",
        ML + Inches(0.15), ty + Inches(2.01),
        lw - Inches(0.25), Inches(1.0), size=10.5, color=TEXT, mono=True)

    txt(s, "MCP SERVERS IN USE TODAY",
        rx, ty - Inches(0.28), rw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    tbl(s, rx, ty, rw, Inches(3.0),
        [
            ("Atlassian", "Jira tickets · Confluence pages · Bitbucket PRs"),
            ("ide",       "Live diagnostics · in-editor code execution"),
            ("Figma",     "Design data · component specs · asset export"),
            ("XcodeBuildMCP", "Build · test · simulate · screenshot on iOS"),
        ],
        headers=["Server", "Reaches into"],
        col0_color=PURPLE, col1_color=MUTED,
        col_widths=[Inches(2.0), Inches(3.5)])

    caption(s, "MCP tools are just tool calls — no special wiring in the agent. "
            "Any agent with the right MCP server configured can reach any external system.",
            ty + Inches(3.15))


def s12_collaboration(prs):
    s = blank(prs)
    lbl(s, "Research Track A — Collaboration")
    hl(s, "One repo. All agents, skills, and references.", size=22)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.45)

    txt(s, "CURRENT APPROACH",
        ML, ty - Inches(0.28), lw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    tree = ("software-dev-agentic/\n"
            "├── lib/core/          ← shared agents + skills\n"
            "├── lib/platforms/     ← web · ios · flutter · android\n"
            "│   ├── web/\n"
            "│   ├── ios/\n"
            "│   ├── flutter/\n"
            "│   └── android/\n"
            "└── docs/principles/   ← reference docs")
    box(s, ML, ty, lw, Inches(3.2), SURFACE, BORDER)
    txt(s, tree, ML + Inches(0.15), ty + Inches(0.15),
        lw - Inches(0.25), Inches(2.9), size=12, color=TEXT, mono=True)
    caption(s, "Single source of truth — one PR, one version, all platforms.",
            ty + Inches(3.3))

    txt(s, "OPEN QUESTIONS",
        rx, ty - Inches(0.28), rw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    qs = [
        ("Contribution workflow", "Who reviews what, across platforms?"),
        ("Ownership model", "Core team vs. platform teams?"),
        ("Versioning", "Per platform or a single version for all?"),
        ("Review process", "How do changes get validated before merge?"),
    ]
    qy = ty
    for title, detail in qs:
        box(s, rx, qy, rw, Inches(0.72), SURFACE, BORDER)
        txt(s, "→  " + title + " — " + detail,
            rx + Inches(0.15), qy + Inches(0.18),
            rw - Inches(0.25), Inches(0.42), size=11.5, color=MUTED)
        qy += Inches(0.82)


def s13_distribution(prs):
    s = blank(prs)
    lbl(s, "Research Track B — Distribution")
    hl(s, "Two paths. Each project picks the one that fits.", size=22)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.45)
    th = Inches(4.3)

    box(s, ML, ty, lw, th, SURFACE, BORDER)
    txt(s, "PATH A — SUBMODULE + SYMLINKS",
        ML + Inches(0.15), ty + Inches(0.1),
        lw - Inches(0.25), Inches(0.22), size=8, bold=True, color=ACCENT)
    code_a = ("project-root/\n"
              "├── software-dev-agentic/  # git submodule\n"
              "└── .claude/\n"
              "    ├── agents → ../software-dev-agentic/lib/\n"
              "    └── skills → ../software-dev-agentic/lib/\n\n"
              "✓ Proven on iOS · Flutter · Android · Web\n"
              "✓ Each project pins a commit\n"
              "✓ Works offline, no build step")
    txt(s, code_a, ML + Inches(0.15), ty + Inches(0.38),
        lw - Inches(0.25), th - Inches(0.5), size=11, color=TEXT, mono=True)

    box(s, rx, ty, rw, th, SURFACE, BORDER)
    txt(s, "PATH B — CLAUDE PLUGIN",
        rx + Inches(0.15), ty + Inches(0.1),
        rw - Inches(0.25), Inches(0.22), size=8, bold=True, color=ACCENT)
    code_b = ("# One command to install\n"
              "curl -fsSL .../install-plugin.sh | bash -s \\\n"
              "  -- --platform=flutter-mobile\n\n"
              "# Auto-patched into .claude/settings.json\n"
              '{\n'
              '  "enabledPlugins": {\n'
              '    "sda-flutter-mobile@sda": true\n'
              '  },\n'
              '  "skillListingBudgetFraction": 0.03\n'
              '}\n\n'
              "✓ No submodule or symlinks\n"
              "✓ Project-scoped — no cross-platform conflicts\n"
              "✓ Auto-updates on every /release")
    txt(s, code_b, rx + Inches(0.15), ty + Inches(0.38),
        rw - Inches(0.25), th - Inches(0.5), size=11, color=TEXT, mono=True)

    caption(s, "Both paths are supported. Submodule fits repos that already manage dependencies via git. "
            "Plugin fits new projects — one command, zero manual wiring.", H - Inches(0.85))


def s_knowledge_arch(prs):
    s = blank(prs)
    lbl(s, "Research Track C — Knowledge Architecture")
    hl(s, "One source of truth. How does it scale?", size=22)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty = MT + Inches(1.45)

    txt(s, "WHERE WE ARE TODAY",
        ML, ty - Inches(0.28), lw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    tree = ("Topic   domain · data · presentation · design-system\n"
            "  │  one reference file per topic per platform\n"
            "  │\n"
            "  ├── -theory.md   what it IS · platform-agnostic\n"
            "  ├── -impl.md     how to build it in this language\n"
            "  └── -catalog.md  queryable symbol inventory\n\n"
            "Retrieval: Grep-first · section offset · never full-file")
    box(s, ML, ty, lw, Inches(2.5), SURFACE, BORDER)
    txt(s, tree, ML + Inches(0.15), ty + Inches(0.15),
        lw - Inches(0.25), Inches(2.3), size=11.5, color=TEXT, mono=True)
    caption(s, "Works well at current scale. The question is what breaks as the system grows.",
            ty + Inches(2.65))

    txt(s, "OPEN QUESTIONS",
        rx, ty - Inches(0.28), rw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    qs = [
        ("Authoring conventions",
         "How do we keep Term headings and file structure consistent as more teams add docs?"),
        ("Versioning",
         "When a platform pattern changes, how does the reference file track old and new "
         "versions without breaking agents that haven't migrated?"),
        ("Retrieval at scale",
         "Grep works today. At what point does a vector index, catalog service, "
         "or structured query replace it?"),
        ("Cross-platform knowledge",
         "Some patterns are shared across platforms. Where does that live, and how do "
         "agents load it without duplicating it?"),
    ]
    qy = ty
    for title, detail in qs:
        bh = Inches(0.9)
        box(s, rx, qy, rw, bh, SURFACE, BORDER)
        txt(s, "→  " + title, rx + Inches(0.15), qy + Inches(0.08),
            rw - Inches(0.25), Inches(0.25), size=11.5, bold=True, color=TEXT)
        txt(s, detail, rx + Inches(0.15), qy + Inches(0.35),
            rw - Inches(0.25), Inches(0.5), size=10.5, color=MUTED)
        qy += bh + Inches(0.12)


def s14_expansion(prs):
    s = blank(prs)
    lbl(s, "Research Track D — Expansion")
    hl(s, "From Implementation to every process.", size=22)

    phases = [
        ("01", "Requirement", "Spec agents, PRD generation, risk analysis", False),
        ("02", "Design",      "Architecture agents, diagram generation", False),
        ("03", "Implementation", "Developer · Debugger · Auditor — active today", True),
        ("04", "Testing",    "Test generation, coverage analysis agents", False),
        ("05", "Delivery",   "Release agents, deploy pipeline automation", False),
    ]
    cw = Inches(2.15)
    gap = Inches(0.1)
    ty = MT + Inches(1.45)
    th = Inches(2.0)
    for i, (num, title, desc, active) in enumerate(phases):
        cx = ML + i * (cw + gap)
        fill = ACCENT_BG if active else SURFACE
        border = ACCENT if active else BORDER
        alpha = 1.0 if active else 0.6
        box(s, cx, ty, cw, th, fill, border)
        txt(s, num, cx + Inches(0.15), ty + Inches(0.12),
            cw - Inches(0.25), Inches(0.22), size=8.5, bold=True, color=ACCENT)
        txt(s, title, cx + Inches(0.15), ty + Inches(0.38),
            cw - Inches(0.25), Inches(0.38), size=13, bold=True, color=TEXT)
        txt(s, desc, cx + Inches(0.15), ty + Inches(0.85),
            cw - Inches(0.25), Inches(0.9), size=10.5, color=MUTED)

    lw = Inches(5.5)
    rw = Inches(5.5)
    rx = ML + lw + Inches(0.33)
    ty2 = ty + th + Inches(0.3)

    caption(s, "No progress yet — runs as a parallel research track. Expansion defines the end state: "
            "specialized agents across every SDLC process, not just Implementation.", ty2)

    txt(s, "OPEN QUESTIONS",
        rx, ty2, rw, Inches(0.22), size=8.5, bold=True, color=MUTED)
    box(s, rx, ty2 + Inches(0.28), rw, Inches(0.8), SURFACE, BORDER)
    txt(s, "→  Orchestrating every process — how does one agent hand off to the next "
        "across the full pipeline?",
        rx + Inches(0.15), ty2 + Inches(0.38),
        rw - Inches(0.25), Inches(0.6), size=11.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

BUILDERS = [
    s01_cover,
    s02_background,
    s03_shift,
    s04_goal,
    s05_vision,
    s06_roadmap,
    s07_anatomy,
    s07b_capability_matrix,
    s_persona,
    s_example,
    s_why_not_e2e,
    s_evolution,
    s_reference_hl,
    s_reference_ll,
    s_skills_hl,
    s_skill_flow,
    s_skills_ll,
    s_agents_hl,
    s_agents_why,
    s_agent_anatomy_ll,
    s_mcp_hl,
    s_mcp_ll,
    s12_collaboration,
    s13_distribution,
    s_knowledge_arch,
    s14_expansion,
]

if __name__ == "__main__":
    out = os.path.join(os.path.dirname(__file__), "agentic-deck.pptx")
    prs = new_prs()
    for i, build in enumerate(BUILDERS, 1):
        print(f"  [{i:02d}/{len(BUILDERS)}] {build.__name__}")
        build(prs)
    prs.save(out)
    print(f"\nSaved → {out}")
    print("Import into Google Slides: File → Import slides → Upload")
